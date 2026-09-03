#!/usr/bin/env python3
"""bb2dash :: extract_text.py
Extract text units from pdf/docx/pptx/xlsx files for bb_file_text.
Usage: extract_text.py <file> [...]  -> prints JSON list of {file, unit_kind, unit_no, text}
"""
import sys, json, subprocess, os
def pdf(p):
    out = subprocess.run(["pdftotext", "-layout", p, "-"], capture_output=True, text=True).stdout
    return [{"unit_kind": "page", "unit_no": i+1, "text": t} for i, t in enumerate(out.split("\f")) if t.strip()]
def docx_(p):
    import docx
    d = docx.Document(p); parts = [x.text for x in d.paragraphs if x.text.strip()]
    for t in d.tables:
        for r in t.rows: parts.append(" | ".join(c.text.strip() for c in r.cells))
    return [{"unit_kind": "doc", "unit_no": 1, "text": "\n".join(parts)}]
def pptx_(p):
    from pptx import Presentation
    units = []
    for i, s in enumerate(Presentation(p).slides, 1):
        parts = []
        for sh in s.shapes:
            if sh.has_text_frame: parts.append(sh.text_frame.text)
            if getattr(sh, "has_table", False) and sh.has_table:
                for r in sh.table.rows: parts.append(" | ".join(c.text for c in r.cells))
        if s.has_notes_slide and s.notes_slide.notes_text_frame.text.strip():
            parts.append("[notes] " + s.notes_slide.notes_text_frame.text)
        t = "\n".join(x for x in parts if x and x.strip())
        if t.strip(): units.append({"unit_kind": "slide", "unit_no": i, "text": t})
    return units
def xlsx_(p):
    import openpyxl
    wb = openpyxl.load_workbook(p, data_only=True, read_only=True); units = []
    for i, ws in enumerate(wb.worksheets, 1):
        rows = [" | ".join("" if v is None else str(v) for v in r) for r in ws.iter_rows(values_only=True)]
        rows = [r for r in rows if r.replace("|", "").strip()]
        if rows: units.append({"unit_kind": "sheet", "unit_no": i, "text": f"# {ws.title}\n" + "\n".join(rows)})
    return units
H = {".pdf": pdf, ".docx": docx_, ".pptx": pptx_, ".xlsx": xlsx_}
out = []
for p in sys.argv[1:]:
    ext = os.path.splitext(p)[1].lower()
    try:
        units = H[ext](p) if ext in H else []
        status = "extracted" if units else ("na" if ext not in H else "failed")
    except Exception as e:
        units, status = [], f"failed: {e}"
    out.append({"file": p, "status": status, "units": units})
print(json.dumps(out))
